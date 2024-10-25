import cv2
import numpy as np
from pptx import Presentation

from HandTracker import HandDetector

# Variables
width, height = 1280, 720
presentation_file = "C:/Users/Bhanu/Documents/21.pptx"  # Change to your presentation file
slide_num = 0
hs, ws = int(120 * 1.2), int(213 * 1.2)
ge_thresh_y = 400
ge_thresh_x = 750
gest_done = False
gest_counter = 0
delay = 5  # Reduced delay for faster response
drawing = False  # Add a flag to track drawing state
annotations = [[]]
annot_num = 0
annot_start = False

# Pointer settings
pointer_color = (0, 0, 255)  # Default color: red
pointer_shape = "circle"      # Default shape: circle

# Load PowerPoint presentation
presentation = Presentation(presentation_file)
slides = list(presentation.slides)

# Camera Setup
cap = cv2.VideoCapture(0)
cap.set(3, width)
cap.set(4, height)

# HandDetector
detector = HandDetector(detectionCon=0.6, maxHands=1)

# Function to change pointer color and shape
def change_pointer_color(color):
    global pointer_color
    pointer_color = color

# Define gestures for changing pointer color
color_gestures = {
    (1, 0, 0, 0, 0): (0, 0, 255),  # Red
    (0, 1, 0, 0, 0): (0, 255, 0),  # Green
    (0, 0, 1, 0, 0): (255, 0, 0),  # Blue
}

# Define a flag to indicate color selection mode
color_selection_mode = False

while True:
    # Get image frame
    success, frame = cap.read()
    frame = cv2.flip(frame, 1)
    slide_current = np.zeros((height, width, 3), np.uint8)  # Blank slide if no PowerPoint slide

    # Load current slide from PowerPoint
    if slide_num < len(slides):
        slide = slides[slide_num]
        for shape in slide.shapes:
            if hasattr(shape, 'image'):
                # Handle image shapes
                image = shape.image
                image_bytes = image.blob
                nparr = np.frombuffer(image_bytes, np.uint8)
                slide_current = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
                slide_current = cv2.resize(slide_current, (width, height))

    # Find the hand and its landmarks
    hands, frame = detector.findHands(frame)

    # Draw Gesture Threshold line
    cv2.line(frame, (width, 0), (ge_thresh_x, ge_thresh_y), (0, 255, 0), 5)

    if hands and gest_done is False:  # If hand is detected
        hand = hands[0]
        cx, cy = hand["center"]
        lm_list = hand["lmList"]  # List of 21 Landmark points
        fingers = detector.fingersUp(hand)

        # Check for color selection gesture (e.g., holding four fingers up)
        if fingers == [1, 1, 1, 1, 0]:
            color_selection_mode = True
        elif fingers != [0, 0, 0, 0, 0]:
            color_selection_mode = False

        # If in color selection mode, check for color change gestures
        if color_selection_mode:
            for gesture, color in color_gestures.items():
                if fingers == gesture:
                    print("Changing pointer color to:", color)  # Debug statement
                    change_pointer_color(color)
                    # Reset color selection mode
                    color_selection_mode = False

        # Constrain values for easier drawing
        x_val = int(np.interp(cx, [width // 2, width], [0, width]))
        y_val = int(np.interp(cy, [150, height - 150], [0, height]))
        index_fing = x_val, y_val

        if cy < ge_thresh_y and cx > ge_thresh_x:
            annot_start = False

            # gest_1 (previous)
            if fingers == [1, 0, 0, 0, 0]:
                annot_start = False
                if slide_num > 0:
                    gest_done = True
                    slide_num -= 1
                    annotations = [[]]
                    annot_num = 0

            # gest_2 (next)
            if fingers == [0, 0, 0, 0, 1]:
                annot_start = False
                if slide_num < len(slides) - 1:
                    gest_done = True
                    slide_num += 1
                    annotations = [[]]
                    annot_num = 0

            # gest_3 (clear screen)
            if fingers == [1, 1, 1, 1, 1]:
                if annotations:
                    annot_start = False
                    if annot_num >= 0:
                        annotations.clear()
                        annot_num = 0
                        gest_done = True
                        annotations = [[]]

        # gest_4 (show pointer)
        if fingers == [0, 1, 1, 0, 0]:
            cv2.circle(slide_current, index_fing, 4, pointer_color, cv2.FILLED)
            annot_start = False
            drawing = False  # Disable drawing when showing pointer

        # Define a scaling factor for sensitivity adjustment
        sensitivity_factor = 0.6  # Adjust this value as needed

        # gest_5 (draw)
        if fingers == [0, 1, 0, 0, 0]:
            if annot_start is False:
                annot_start = True
                annot_num += 1
                annotations.append([])
                drawing = True  # Enable drawing when gesture detected
                prev_index_fing = index_fing  # Store the initial finger position
            if drawing:  # Only draw when the gesture is active
                annotations[annot_num].append(index_fing)
                cv2.circle(slide_current, index_fing, 4, pointer_color, cv2.FILLED)
                # Adjust sensitivity by multiplying the difference by the scaling factor
                delta_x = int((index_fing[0] - prev_index_fing[0]) * sensitivity_factor)
                delta_y = int((index_fing[1] - prev_index_fing[1]) * sensitivity_factor)
                # Draw a line from the previous finger position to the current finger position
                cv2.line(slide_current, prev_index_fing, (prev_index_fing[0] + delta_x, prev_index_fing[1] + delta_y),
                         pointer_color, 6)
                prev_index_fing = (
                prev_index_fing[0] + delta_x, prev_index_fing[1] + delta_y)  # Update the previous finger position

        # gest_6 (erase)
        if fingers == [0, 1, 1, 1, 0]:
            if annotations:
                annot_start = False
                if annot_num >= 0:
                    annotations.pop(-1)
                    annot_num -= 1
                    gest_done = True

    else:
        annot_start = False
    gesture_delay = 15
    # Gesture Performed Iterations:
    if gest_done:
        gest_counter += 1
        if gest_counter > delay:
            gest_counter = 0
            gest_done = False

    for annotation in annotations:
        if len(annotation) > 1:
            for i in range(1, len(annotation)):
                cv2.line(slide_current, annotation[i - 1], annotation[i], pointer_color, 6)

    # Adding cam img on slides
    composite_img = np.zeros((max(height, hs), width + ws, 3), dtype=np.uint8)

    # Define the size of the camera box
    camera_box_width = ws
    camera_box_height = hs

    # Calculate the width of the combined area (presentation slide + camera box)
    combined_width = width + camera_box_width

    # Calculate the height of the combined area (presentation slide or camera box, whichever is taller)
    combined_height = max(height, camera_box_height)

    # Calculate the offset to center the camera box vertically
    camera_offset_y = (combined_height - camera_box_height) // 2

    # Calculate the positions for the presentation slide and camera screen
    slide_top_left = (0, 0)
    camera_bottom_right = (combined_width, combined_height)

    # Add the presentation slide to the top left, cropping to the right if needed
    composite_img[slide_top_left[1]:slide_top_left[1] + height,
    slide_top_left[0]:slide_top_left[0] + width] = slide_current

    # Resize the camera frame
    resized_frame = cv2.resize(frame, (camera_box_width, camera_box_height))

    # Add the camera screen to the bottom right as an overlay
    composite_img[camera_offset_y:camera_offset_y + camera_box_height,
    width:combined_width] = resized_frame

    cv2.imshow("Slides", composite_img)

    key = cv2.waitKey(1)
    if key == ord('q'):
        break

# Release the camera and close OpenCV windows
cap.release()
cv2.destroyAllWindows()
